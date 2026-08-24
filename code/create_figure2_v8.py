"""
Figure 2 v8 (revised): Two-Stage Moderated Mediation Model
自动读取最新 STEP4 和 STEP5 输出：
- output/STEP5_wald_tests_complete.csv (Wald χ² tests)
- output/STEP4_new_wald_tests.csv (a-path Wald tests, indirect-effect Wald)
- output/STEP4_new_direct_paths_wald.csv (c' path Wald tests)
- output/STEP5_moderated_mediation_BC_results.csv (b path)

修订布局：调整节点位置和标签间距，避免重叠。
"""
import os
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from matplotlib import rcParams

OUT_DIR = r"d:\Study on Youth Suicide\output"

# ============================================================
# 1. 读取数据
# ============================================================
# a-path Wald tests (使用 indirect-effect Wald tests，与 figure_combined.tex 一致)
wald4 = pd.read_csv(os.path.join(OUT_DIR, "STEP4_new_wald_tests.csv"))
def get_wald4(row_name):
    row = wald4[wald4.iloc[:, 0] == row_name]
    if row.empty:
        return None
    chi2 = float(row['chi2'].values[0])
    p = float(row['pvalue'].values[0])
    return {'chi2': chi2, 'p': p}

# Figure 2: X_1=Bullying, X_2=School, X_3=Internet
a1 = get_wald4("Bullying_Victimization")
a2 = get_wald4("School_Avoidance")
a3 = get_wald4("Internet_Dependency")

# c' path Wald tests (from STEP4_new_direct_paths_wald.csv)
wald4c = pd.read_csv(os.path.join(OUT_DIR, "STEP4_new_direct_paths_wald.csv"))
def get_wald4c(row_name):
    row = wald4c[wald4c.iloc[:, 0] == row_name]
    if row.empty:
        return None
    chi2 = float(row['chi2'].values[0])
    p = float(row['pvalue'].values[0])
    return {'chi2': chi2, 'p': p}

c1 = get_wald4c("Injury_cprime")
c2 = get_wald4c("School_cprime")
c3 = get_wald4c("Internet_cprime")

# STEP5 moderation Wald tests
wald5 = pd.read_csv(os.path.join(OUT_DIR, "STEP5_wald_tests_complete.csv"))
def get_wald5(test_name):
    row = wald5[wald5['test'] == test_name]
    if row.empty:
        return None
    chi2 = float(row['chi2'].values[0])
    p = float(row['pvalue'].values[0])
    df = int(row['df'].values[0])
    return {'chi2': chi2, 'p': p, 'df': df}

school_a    = get_wald5("School_XtoM_moderation")
school_c    = get_wald5("School_XtoY_moderation")
school_t    = get_wald5("School_total_moderation")
injury_a    = get_wald5("Injury_XtoM_moderation")
injury_c    = get_wald5("Injury_XtoY_moderation")
injury_t    = get_wald5("Injury_total_moderation")
inet_a      = get_wald5("Internet_XtoM_moderation")
inet_c      = get_wald5("Internet_XtoY_moderation")
inet_t      = get_wald5("Internet_total_moderation")
bz_wald     = get_wald5("MtoY_moderation_bz")

# b path (from STEP5)
bc_df = pd.read_csv(os.path.join(OUT_DIR, "STEP5_moderated_mediation_BC_results.csv"))
def get_bc(label):
    row = bc_df[bc_df['label'] == label]
    if row.empty:
        return None
    est = float(row['est'].values[0])
    lo = float(row['ci_lower'].values[0])
    hi = float(row['ci_upper'].values[0])
    sig = lo > 0 or hi < 0
    return {'est': est, 'lo': lo, 'hi': hi, 'sig': sig}

b_path = get_bc("b")

# ============================================================
# 2. 格式化辅助函数
# ============================================================
def fmt_chi2_label(d, prefix, df_val=3):
    if d is None: return ""
    p = d['p']
    if p < .001: sig = "***"
    elif p < .01: sig = "**"
    elif p < .05: sig = "*"
    elif p < .10: sig = "†"
    else: sig = ""
    return f"{prefix}: χ²({df_val}) = {d['chi2']:.2f}{sig}"

def fmt_b_path(b):
    if b is None: return ""
    star = "*" if b['sig'] else ""
    return f"b = {b['est']:.3f}{star} [{b['lo']:.2f}, {b['hi']:.2f}]"

# ============================================================
# 3. 设置 matplotlib
# ============================================================
rcParams['font.family'] = 'sans-serif'
rcParams['font.sans-serif'] = ['Arial', 'DejaVu Sans']
rcParams['mathtext.fontset'] = 'custom'
rcParams['mathtext.rm'] = 'Arial'
rcParams['mathtext.it'] = 'Arial:italic'
rcParams['mathtext.bf'] = 'Arial:bold'

# 更大的画布以容纳更多标签
fig, ax = plt.subplots(1, 1, figsize=(17, 12))
ax.set_xlim(-8, 11)
ax.set_ylim(-8, 6)
ax.set_aspect('equal')
ax.axis('off')

def draw_box(ax, x, y, text, width=2.8, height=0.95):
    box = FancyBboxPatch((x - width/2, y - height/2), width, height,
                          boxstyle="round,pad=0.12", fc='white', ec='black', lw=0.8)
    ax.add_patch(box)
    ax.text(x, y, text, ha='center', va='center', fontsize=10,
            fontfamily='Arial', linespacing=1.3)

def draw_arrow(ax, start, end, style='sig', label='', label_pos=0.5,
               label_offset=(0, 0), curvature=0, zorder=3, label_fontsize=7):
    if style == 'sig':          color, lw, ls = '#333333', 2.0, '-'
    elif style == 'ns':         color, lw, ls = '#888888', 1.3, '--'
    elif style == 'direct':     color, lw, ls = '#666666', 0.9, '-'
    elif style == 'mod_sig':    color, lw, ls = '#555555', 1.1, '-'
    else:                       color, lw, ls = '#999999', 1.1, '--'

    if abs(curvature) > 0.01:
        arrow = FancyArrowPatch(start, end,
                                connectionstyle=f"arc3,rad={curvature}",
                                arrowstyle='-|>', mutation_scale=14,
                                color=color, lw=lw, linestyle=ls, zorder=zorder)
    else:
        arrow = FancyArrowPatch(start, end,
                                arrowstyle='-|>', mutation_scale=14,
                                color=color, lw=lw, linestyle=ls, zorder=zorder)
    ax.add_patch(arrow)

    if label:
        if abs(curvature) > 0.01:
            mid_x = (start[0] + end[0]) / 2 + label_offset[0]
            mid_y = (start[1] + end[1]) / 2 + curvature * 1.8 + label_offset[1]
        else:
            mid_x = start[0] + (end[0] - start[0]) * label_pos + label_offset[0]
            mid_y = start[1] + (end[1] - start[1]) * label_pos + label_offset[1]
        ax.text(mid_x, mid_y, label, ha='center', va='center',
                fontsize=label_fontsize, fontfamily='Arial',
                bbox=dict(boxstyle='round,pad=0.18', fc='white', ec='none', alpha=0.92),
                zorder=10, linespacing=1.25)

# ============================================================
# 4. 节点位置（更分散以避免重叠）
# ============================================================
AD  = (-5, 3)
BV  = (-5, 0)
IRI = (-5, -3)
M   = (1.5, 0)
Y   = (7.5, 0)
Z   = (1.5, -5.5)

hw, hh = 1.4, 0.47

draw_box(ax, AD[0], AD[1], 'Academic-disengagement\nprofiles ($X_2$)')
draw_box(ax, BV[0], BV[1], 'Bullying-victimization\nprofiles ($X_1$)')
draw_box(ax, IRI[0], IRI[1], 'Internet-related-impairment\nprofiles ($X_3$)')
draw_box(ax, M[0], M[1], 'Suicidal thoughts /\npreparation ($M$)')
draw_box(ax, Y[0], Y[1], 'Self-harm ($Y$)')
draw_box(ax, Z[0], Z[1], 'Perceived family\nconnectedness ($Z$)')

# ============================================================
# 5. a paths: X -> M (Wald)
# ============================================================
# AD → M
label_ad = fmt_chi2_label(a2, "$a_2$")
draw_arrow(ax, (AD[0]+hw, AD[1]), (M[0]-hw, M[1]+hh*0.8),
           style='sig' if (a2 and a2['p']<.05) else 'ns',
           label=label_ad,
           label_offset=(0, 0.30), label_fontsize=7.5)

# BV → M
label_bv = fmt_chi2_label(a1, "$a_1$")
draw_arrow(ax, (BV[0]+hw, BV[1]), (M[0]-hw, M[1]),
           style='sig' if (a1 and a1['p']<.05) else 'ns',
           label=label_bv,
           label_offset=(0, 0.30), label_fontsize=7.5)

# IRI → M
label_iri = fmt_chi2_label(a3, "$a_3$")
draw_arrow(ax, (IRI[0]+hw, IRI[1]), (M[0]-hw, M[1]-hh*0.8),
           style='sig' if (a3 and a3['p']<.05) else 'ns',
           label=label_iri,
           label_offset=(0, -0.30), label_fontsize=7.5)

# ============================================================
# 6. b path: M → Y
# ============================================================
b_label = fmt_b_path(b_path)
draw_arrow(ax, (M[0]+hw, M[1]), (Y[0]-hw, Y[1]),
           style='sig',
           label=f"b: {b_label}" if b_label else "b",
           label_offset=(0, 0.32), label_fontsize=7.5)

# ============================================================
# 7. Direct effects X → Y (c')
# 使用 c' path Wald (χ²(3)) from STEP4_new_direct_paths_wald.csv
# ============================================================
# AD → Y direct
c2_label = fmt_chi2_label(c2, "c'$_2$")
draw_arrow(ax, (AD[0]+hw, AD[1]-hh*0.6), (Y[0]-hw, Y[1]+hh*0.7),
           style='direct',
           label=c2_label,
           label_pos=0.72, label_offset=(0, 0.20), label_fontsize=7)

# BV → Y direct (curved)
c1_label = fmt_chi2_label(c1, "c'$_1$")
draw_arrow(ax, (BV[0]+hw, BV[1]), (Y[0]-hw, Y[1]+hh*0.7),
           style='direct', curvature=+0.42,
           label=c1_label,
           label_offset=(0, 0.36), label_fontsize=7)

# IRI → Y direct
c3_label = fmt_chi2_label(c3, "c'$_3$")
draw_arrow(ax, (IRI[0]+hw, IRI[1]+hh*0.6), (Y[0]-hw, Y[1]-hh*0.7),
           style='direct',
           label=c3_label,
           label_pos=0.72, label_offset=(0, -0.20), label_fontsize=7)

# ============================================================
# 8. Stage 1 Moderation: Z → a-path midpoints
# ============================================================
# Z → AD→M: a_2z
if school_a:
    style = 'mod_sig' if school_a['p'] < .05 else 'mod_ns'
    label = fmt_chi2_label(school_a, "$a_{2z}$")
    ad_m_mid = ((AD[0]+hw + M[0]-hw)/2 - 0.5, (AD[1] + M[1]+hh*0.8)/2)
    draw_arrow(ax, (Z[0]-hh-0.3, Z[1]+hh*0.9), ad_m_mid,
               style=style, label=label,
               label_pos=0.30, label_offset=(-0.05, 0.0), zorder=4, label_fontsize=7)

# Z → BV→M: a_1z
if injury_a:
    style = 'mod_sig' if injury_a['p'] < .05 else 'mod_ns'
    label = fmt_chi2_label(injury_a, "$a_{1z}$")
    bv_m_mid = ((BV[0]+hw + M[0]-hw)/2, (BV[1] + M[1])/2)
    draw_arrow(ax, (Z[0]-0.5, Z[1]+hh*0.9), bv_m_mid,
               style=style, label=label,
               label_pos=0.30, label_offset=(-0.05, 0.0), zorder=4, label_fontsize=7)

# Z → IRI→M: a_3z
if inet_a:
    style = 'mod_sig' if inet_a['p'] < .05 else 'mod_ns'
    label = fmt_chi2_label(inet_a, "$a_{3z}$")
    iri_m_mid = ((IRI[0]+hw + M[0]-hw)/2 + 0.5, (IRI[1] + M[1]-hh*0.8)/2)
    draw_arrow(ax, (Z[0]+0.3, Z[1]+hh*0.9), iri_m_mid,
               style=style, label=label,
               label_pos=0.30, label_offset=(0.05, 0.0), zorder=4, label_fontsize=7)

# ============================================================
# 9. Stage 2 Moderation: Z → b-path midpoint
# ============================================================
if bz_wald:
    style = 'mod_sig' if bz_wald['p'] < .05 else 'mod_ns'
    label = fmt_chi2_label(bz_wald, "$b_z$", df_val=1)
    m_y_mid = ((M[0]+hw + Y[0]-hw)/2, (M[1] + Y[1])/2 + 0.6)
    draw_arrow(ax, (Z[0]+hh+0.5, Z[1]+hh*0.9), m_y_mid,
               style=style, label=label,
               label_pos=0.40, label_offset=(0.1, 0.0), zorder=4, label_fontsize=7)

# ============================================================
# 10. c' path moderation: Z → direct effect midpoints
# ============================================================
# Z → AD→Y: c_2z
if school_c:
    style = 'mod_sig' if school_c['p'] < .05 else 'mod_ns'
    label = fmt_chi2_label(school_c, "$c_{2z}$")
    ad_y_mid = ((AD[0]+hw + Y[0]-hw)*0.50, (AD[1]-hh*0.6 + Y[1]+hh*0.7)*0.50)
    draw_arrow(ax, (Z[0]-1.0, Z[1]+hh*0.9), ad_y_mid,
               style=style, label=label,
               label_pos=0.40, label_offset=(-0.05, 0.0), zorder=4, label_fontsize=7)

# Z → BV→Y: c_1z
if injury_c:
    style = 'mod_sig' if injury_c['p'] < .05 else 'mod_ns'
    label = fmt_chi2_label(injury_c, "$c_{1z}$")
    bv_y_mid = ((BV[0]+hw + Y[0]-hw)*0.55, (BV[1] + Y[1]+hh*0.7)*0.55)
    draw_arrow(ax, (Z[0]+0.0, Z[1]+hh*0.9), bv_y_mid,
               style=style, label=label,
               label_pos=0.40, label_offset=(0.05, 0.0), zorder=4, label_fontsize=7)

# Z → IRI→Y: c_3z
if inet_c:
    style = 'mod_sig' if inet_c['p'] < .05 else 'mod_ns'
    label = fmt_chi2_label(inet_c, "$c_{3z}$")
    iri_y_mid = ((IRI[0]+hw + Y[0]-hw)*0.50, (IRI[1]+hh*0.6 + Y[1]-hh*0.7)*0.50)
    draw_arrow(ax, (Z[0]+1.0, Z[1]+hh*0.9), iri_y_mid,
               style=style, label=label,
               label_pos=0.40, label_offset=(0.05, 0.0), zorder=4, label_fontsize=7)

# ============================================================
# 11. Legend
# ============================================================
legend_items = [
    ('sig',     'Significant $a/b$ path'),
    ('ns',      'Non-significant path'),
    ('direct',  'Direct effect ($X \\to Y$)'),
    ('mod_sig', 'Significant moderation'),
    ('mod_ns',  'Non-significant moderation'),
]
lx, ly = 7.5, -7.5
for i, (style, text) in enumerate(legend_items):
    yp = ly + i * 0.40
    if style == 'sig':     ax.plot([lx, lx+0.7], [yp, yp], color='#333333', lw=2.0)
    elif style == 'ns':    ax.plot([lx, lx+0.7], [yp, yp], color='#888888', lw=1.3, ls='--')
    elif style == 'direct':ax.plot([lx, lx+0.7], [yp, yp], color='#666666', lw=0.9)
    elif style == 'mod_sig':ax.plot([lx, lx+0.7], [yp, yp], color='#555555', lw=1.1)
    else:                  ax.plot([lx, lx+0.7], [yp, yp], color='#999999', lw=1.1, ls='--')
    ax.text(lx + 0.85, yp, text, fontsize=8, fontfamily='Arial', va='center', ha='left')

# ============================================================
# 12. Title and note
# ============================================================
ax.text(2.5, 5.5,
        'Figure 2. Two-Stage Moderated Mediation Model',
        ha='center', va='center', fontsize=14, fontweight='bold', fontfamily='Arial')

note_text = (
    'Note. AD = academic disengagement; BV = bullying victimization; '
    'IRI = internet-related impairment; M = suicidal thoughts/preparation; '
    'Y = self-harm; Z = perceived family connectedness.\n'
    '$\\chi^2$ = Wald chi-square test; $b$ = unstandardized coefficient (BC bootstrap 95% CI). '
    'Moderation arrows point to the midpoint of the moderated path.\n'
    'All coefficients are unstandardized with bias-corrected 95% bootstrap CIs (10,000 resamples). '
    '$^* p < .05$, $^{**} p < .01$, $^{***} p < .001$.'
)
ax.text(Z[0], Z[1] - 2.6, note_text,
        ha='center', va='top', fontsize=8, fontfamily='Arial',
        bbox=dict(boxstyle='round,pad=0.5', fc='#fafafa', ec='#cccccc', lw=0.5),
        linespacing=1.5)

plt.tight_layout()
plt.savefig(os.path.join(OUT_DIR, "Figure2_path_diagram_v8.png"),
            dpi=300, bbox_inches='tight', facecolor='white', edgecolor='none')
print(f"Saved PNG: {os.path.join(OUT_DIR, 'Figure2_path_diagram_v8.png')}")

# 也保存为 PPTX
try:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(8.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img_path = os.path.join(OUT_DIR, "Figure2_path_diagram_v8.png")
    slide.shapes.add_picture(img_path, Inches(0.4), Inches(0.3),
                             width=Inches(12.53), height=Inches(7.9))
    pptx_path = os.path.join(OUT_DIR, "Figure2_path_diagram_v8.pptx")
    prs.save(pptx_path)
    print(f"Saved PPTX: {pptx_path}")
except Exception as e:
    print(f"PPTX save skipped: {e}")